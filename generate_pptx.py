"""
generate_pptx.py - Genera presentacion PowerPoint para Task 2
Metodologia CRISP-ML aplicada al proyecto de British Airways
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

REPORTS_DIR = 'docs'
IMAGES_DIR = 'docs'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x01, 0x1E, 0x41)
BLUE = RGBColor(0x07, 0x5A, 0xAA)
RED = RGBColor(0xE3, 0x18, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF7, 0xF9, 0xFC)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GOLD = RGBColor(0xC8, 0xA9, 0x6E)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
    return txBox

# ============================================================
# SLIDE 1: PORTADA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_shape(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), RED)
add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.05), RED)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             'BRITISH AIRWAYS', font_size=48, bold=True, color=WHITE,
             alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.8),
             'Task 2: Predictive Modeling - Customer Booking Completion',
             font_size=28, bold=False, color=GOLD, alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
             'Programa de Simulacion de Data Science | Forage',
             font_size=18, color=RGBColor(0x9C, 0xA3, 0xAF), alignment=PP_ALIGN.LEFT)

add_bullet_slide(slide, Inches(1), Inches(4.4), Inches(11), Inches(2.5), [
    'Metodologia: CRISP-ML (Cross-Industry Standard Process for Machine Learning)',
    'Problema: Clasificacion binaria para predecir reservas completadas',
    'Dataset: 50,000 registros | 14 variables | 15% tasa de completado',
    'Modelos: Logistic Regression, Random Forest, Gradient Boosting + SMOTE',
], font_size=16, color=RGBColor(0xD1, 0xD5, 0xDB))

add_text_box(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.5),
             'Desarrollado por: Feibert Guzman',
             font_size=14, color=RGBColor(0x6B, 0x72, 0x80), alignment=PP_ALIGN.LEFT)

# ============================================================
# SLIDE 2: BUSINESS UNDERSTANDING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), NAVY)
add_shape(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), RED)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             'Fase 1: Business Understanding', font_size=30, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             'Objetivo del Negocio y del Machine Learning', font_size=20, bold=True, color=NAVY)

add_shape(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), LIGHT)
add_text_box(slide, Inches(1), Inches(2.5), Inches(5), Inches(0.5),
             'Objetivo de Negocio', font_size=20, bold=True, color=NAVY)
add_bullet_slide(slide, Inches(1), Inches(3.1), Inches(5), Inches(3.5), [
    'British Airways quiere entender que',
    'factores influyen en la decision de los',
    'clientes de completar una reserva.',
    '',
    'Identificar patrones de comportamiento',
    'para optimizar estrategias de marketing',
    'y mejorar la experiencia de usuario.',
    '',
    'Reducir la tasa de abandono en el',
    'proceso de reserva online.',
], font_size=15, color=DARK_TEXT)

add_shape(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(4.5), NAVY)
add_text_box(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(0.5),
             'Objetivo de Machine Learning', font_size=20, bold=True, color=GOLD)
add_bullet_slide(slide, Inches(7.2), Inches(3.1), Inches(5), Inches(3.5), [
    'Construir un modelo de clasificacion',
    'binaria para predecir si un cliente',
    'completara la reserva.',
    '',
    'Variable objetivo: booking_complete',
    '  1 = Completo  |  0 = No completo',
    '',
    'Metricas de exito:',
    '  AUC-ROC, F1-Score, Precision, Recall',
], font_size=15, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
             'CRISP-ML Fase 1/6', font_size=12, color=GRAY)

# ============================================================
# SLIDE 3: DATA UNDERSTANDING - EDA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), NAVY)
add_shape(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), RED)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             'Fase 2: Data Understanding - EDA', font_size=30, bold=True, color=WHITE)

stats = [
    ('50,000', 'Registros'),
    ('14', 'Variables'),
    ('14.96%', 'Tasa Completado'),
    ('5.69x', 'Desbalanceo'),
]
for i, (val, label) in enumerate(stats):
    left = Inches(0.8 + i * 3.1)
    shape = add_shape(slide, left, Inches(1.5), Inches(2.8), Inches(1.2), LIGHT)
    add_text_box(slide, left, Inches(1.55), Inches(2.8), Inches(0.6),
                 val, font_size=32, bold=True, color=RED, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.1), Inches(2.8), Inches(0.5),
                 label, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Images grid
img_files = [
    ('target_distribution.png', 'Distribucion de Variable Objetivo'),
    ('categorical_analysis.png', 'Variables Categoricas vs Target'),
    ('correlation_matrix.png', 'Matriz de Correlacion'),
    ('services_analysis.png', 'Analisis de Servicios'),
]
for i, (img, caption) in enumerate(img_files):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(3.0 + row * 2.1)
    img_path = os.path.join(IMAGES_DIR, img)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, Inches(5.8), Inches(1.8))
    add_text_box(slide, left, Inches(top.inches + 1.85), Inches(5.8), Inches(0.3),
                 caption, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
             'CRISP-ML Fase 2/6', font_size=12, color=GRAY)

# ============================================================
# SLIDE 4: DATA PREPARATION - ETL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), NAVY)
add_shape(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), RED)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             'Fase 3: Data Preparation - Pipeline ETL', font_size=30, bold=True, color=WHITE)

phases = [
    ('EXTRACCION', BLUE, [
        'Lectura de customer_booking.csv',
        'Codificacion ISO-8859-1',
        'Validacion de tipos de datos',
        '50,000 registros cargados',
    ]),
    ('TRANSFORMACION', RED, [
        'One-hot encoding categorico',
        'Binning (purchase_lead, stay)',
        'Feature engineering',
        'Manejo de outliers (clipping)',
        '36 features generadas',
    ]),
    ('CARGA', NAVY, [
        'processed_booking_data.csv',
        'encoded_booking_data.csv',
        'Datos listos para modelado',
        'Train-Test split (80/20)',
    ]),
]

for i, (title, color, items) in enumerate(phases):
    left = Inches(0.8 + i * 4.1)
    add_shape(slide, left, Inches(1.5), Inches(3.8), Inches(5.5), LIGHT)
    add_shape(slide, left, Inches(1.5), Inches(3.8), Inches(0.6), color)
    add_text_box(slide, left, Inches(1.55), Inches(3.8), Inches(0.5),
                 f'{" " * 3}{i+1}. {title}', font_size=18, bold=True, color=WHITE)
    add_bullet_slide(slide, left + Inches(0.3), Inches(2.3), Inches(3.2), Inches(4.5),
                     items, font_size=14, color=DARK_TEXT)

add_text_box(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
             'CRISP-ML Fase 3/6  |  etl_pipeline.py', font_size=12, color=GRAY)

# ============================================================
# SLIDE 5: MODELING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), NAVY)
add_shape(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), RED)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             'Fase 4: Modeling - 3 Modelos Entrenados', font_size=30, bold=True, color=WHITE)

models_info = [
    ('Logistic Regression\n+ SMOTE', BLUE, 'Modelo Lineal', [
        'Accuracy: 0.725',
        'AUC-ROC: 0.710',
        'Benchmark basico',
    ]),
    ('Random Forest\n+ SMOTE', RED, 'MEJOR MODELO', [
        'Accuracy: 0.802',
        'AUC-ROC: 0.785',
        'Feature importance nativa',
    ]),
    ('Gradient Boosting\n+ SMOTE', NAVY, 'Ensemble', [
        'Accuracy: 0.785',
        'AUC-ROC: 0.762',
        'Arboles secuenciales',
    ]),
]

for i, (name, color, subtitle, metrics) in enumerate(models_info):
    left = Inches(0.8 + i * 4.1)
    add_shape(slide, left, Inches(1.5), Inches(3.8), Inches(5.5), LIGHT)
    add_shape(slide, left, Inches(1.5), Inches(3.8), Inches(2.2), color)
    add_text_box(slide, left, Inches(1.7), Inches(3.8), Inches(0.8),
                 name, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.5), Inches(3.8), Inches(0.4),
                 subtitle, font_size=14, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_bullet_slide(slide, left + Inches(0.3), Inches(3.9), Inches(3.2), Inches(3.0),
                     metrics, font_size=15, color=DARK_TEXT)

add_text_box(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
             'CRISP-ML Fase 4/6  |  SMOTE usado para balanceo de clases', font_size=12, color=GRAY)

# ============================================================
# SLIDE 6: EVALUATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), NAVY)
add_shape(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), RED)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             'Fase 5: Evaluation - Resultados', font_size=30, bold=True, color=WHITE)

# Results table
table_data = [
    ['Modelo', 'Accuracy', 'Precision', 'Recall', 'F1-Score', 'AUC-ROC'],
    ['Logistic Regression + SMOTE', '0.725', '0.285', '0.610', '0.388', '0.710'],
    ['Random Forest + SMOTE', '0.802', '0.420', '0.520', '0.465', '0.785'],
    ['Gradient Boosting + SMOTE', '0.785', '0.380', '0.490', '0.428', '0.762'],
]

rows = len(table_data)
cols = len(table_data[0])
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.0))
table = table_shape.table

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.name = 'Calibri'
            paragraph.alignment = PP_ALIGN.CENTER
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            elif r == 2:
                paragraph.font.bold = True
                paragraph.font.color.rgb = RED
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        elif r == 2:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# Images
img_path = os.path.join(IMAGES_DIR, 'model_comparison.png')
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.8), Inches(3.8), Inches(5.8), Inches(3.2))

img_path2 = os.path.join(IMAGES_DIR, 'feature_importance.png')
if os.path.exists(img_path2):
    slide.shapes.add_picture(img_path2, Inches(7), Inches(3.8), Inches(5.8), Inches(3.2))

add_text_box(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
             'CRISP-ML Fase 5/6  |  Mejor modelo: Random Forest + SMOTE (AUC-ROC = 0.785)',
             font_size=12, color=GRAY)

# ============================================================
# SLIDE 7: DEPLOYMENT - CONCLUSIONES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), NAVY)
add_shape(slide, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06), RED)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             'Fase 6: Deployment & Conclusiones', font_size=30, bold=True, color=WHITE)

# Top Features
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.5), LIGHT)
add_text_box(slide, Inches(1), Inches(1.7), Inches(5.5), Inches(0.5),
             'Top 5 Features mas Importantes', font_size=20, bold=True, color=NAVY)

features = [
    ('1. purchase_lead', '18.2%', 'Tiempo entre compra y viaje'),
    ('2. wants_extra_baggage', '9.5%', 'Solicitud de equipaje extra'),
    ('3. flight_duration', '7.8%', 'Duracion del vuelo'),
    ('4. booking_origin', '6.2%', 'Pais de origen'),
    ('5. wants_preferred_seat', '5.1%', 'Asiento preferencial'),
]
for i, (feat, pct, desc) in enumerate(features):
    y = Inches(2.4 + i * 0.9)
    add_shape(slide, Inches(1), y, Inches(0.06), Inches(0.7), BLUE)
    add_text_box(slide, Inches(1.2), y, Inches(2), Inches(0.4),
                 feat, font_size=15, bold=True, color=NAVY)
    add_text_box(slide, Inches(3.2), y, Inches(1), Inches(0.4),
                 pct, font_size=15, bold=True, color=RED)
    add_text_box(slide, Inches(1.2), y + Inches(0.35), Inches(5), Inches(0.3),
                 desc, font_size=12, color=GRAY)

# Conclusiones
add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.5), NAVY)
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.5),
             'Conclusiones Clave', font_size=20, bold=True, color=GOLD)
add_bullet_slide(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(1.5), [
    'Desbalanceo de clases manejado con SMOTE',
    'Random Forest obtuvo mejor rendimiento',
    'purchase_lead es el factor mas determinante',
    'Servicios adicionales aumentan conversion',
], font_size=14, color=WHITE)

# Recomendaciones
add_shape(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(2.7), BLUE)
add_text_box(slide, Inches(7.2), Inches(4.5), Inches(5), Inches(0.5),
             'Recomendaciones de Negocio', font_size=20, bold=True, color=WHITE)
add_bullet_slide(slide, Inches(7.2), Inches(5.1), Inches(5), Inches(1.7), [
    'Retargeting para clientes con >30% probabilidad',
    'Mejorar UX en canal mobile',
    'Ofrecer incentivos (equipaje, asientos)',
    'Alertas para reservas con largo purchase_lead',
], font_size=14, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
             'CRISP-ML Fase 6/6  |  Entregables: app.py, index.html, notebook, reporte',
             font_size=12, color=GRAY)

# ============================================================
# SLIDE 8: ENTREGABLES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), RED)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'ENTREGABLES DEL PROYECTO', font_size=36, bold=True, color=WHITE)
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
             'Ciclo de Vida Completo del Machine Learning - CRISP-ML', font_size=18, color=GOLD)

deliverables = [
    ('etl_pipeline.py', 'Pipeline ETL'),
    ('customer_booking_analysis.ipynb', 'Notebook Predictivo'),
    ('app.py', 'Dashboard Streamlit'),
    ('index.html', 'Landing Page'),
    ('README.md', 'Documentacion + Badges'),
    ('docs/reporte_integral.html', 'Reporte Integral'),
]

for i, (file_name, desc) in enumerate(deliverables):
    row = i // 3
    col = i % 3
    left = Inches(1 + col * 4)
    top = Inches(2.2 + row * 2.4)
    add_shape(slide, left, top, Inches(3.5), Inches(1.8),
              RGBColor(0x0A, 0x2A, 0x52))
    add_shape(slide, left, top, Inches(3.5), Inches(0.06), RED)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.3), Inches(3), Inches(0.5),
                 file_name, font_size=16, bold=True, color=GOLD)
    add_text_box(slide, left + Inches(0.3), top + Inches(1), Inches(3), Inches(0.5),
                 desc, font_size=14, color=WHITE)

add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
             'Desarrollado por: Feibert Guzman  |  British Airways Forage Data Science Simulation',
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: GRACIAS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06), RED)

add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.2),
             'GRACIAS', font_size=60, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
             'British Airways - Task 2: Predictive Modeling',
             font_size=24, color=GOLD, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
             'CRISP-ML | ETL | EDA | Machine Learning | Deployment',
             font_size=18, color=RGBColor(0x9C, 0xA3, 0xAF), alignment=PP_ALIGN.CENTER)

# SAVE
output_path = 'docs/British_Airways_Task2_Predictive_Modeling.pptx'
prs.save(output_path)
print(f"Presentacion guardada: {output_path}")
print(f"Total de diapositivas: {len(prs.slides)}")
